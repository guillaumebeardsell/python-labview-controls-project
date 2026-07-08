"""Build the "what's next" team deck as an editable 16:9 PowerPoint.

Companion to build_progress_pptx.py (accomplishments since Jul 2); same palette
and visual language, five slides on the upcoming tasks and objectives:

  1. The road ahead — authority ladder B -> C -> D -> E, gated at every step
  2. B3 — the one LabVIEW build between shadow and command
  3. B4 — the seven failure drills that gate any authority
  4. Phase C — Python commands the bench
  5. What the team supplies + where this lands (commissioning)

    python team-update/build_next_pptx.py  ->  team-update/migration-next-steps.pptx
"""

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = pathlib.Path(__file__).parent

# palette (same as build_pptx.py / build_progress_pptx.py)
INK = RGBColor(0x16, 0x22, 0x2D)
MUTED = RGBColor(0x5C, 0x6B, 0x78)
HAIR = RGBColor(0xE2, 0xE7, 0xEC)
PY = RGBColor(0x5B, 0x46, 0xC9)
PY_BG = RGBColor(0xED, 0xEA, 0xFA)
PY_LN = RGBColor(0xD6, 0xCE, 0xF4)
LV = RGBColor(0x0E, 0x76, 0x7E)
LV_BG = RGBColor(0xE0, 0xF0, 0xF0)
LV_LN = RGBColor(0xC4, 0xE2, 0xE2)
OK = RGBColor(0x1C, 0x91, 0x4F)
OK_BG = RGBColor(0xE3, 0xF4, 0xEB)
OK_LN = RGBColor(0xBF, 0xE6, 0xCE)
WARN = RGBColor(0xB4, 0x5B, 0x0E)
WARN_BG = RGBColor(0xFB, 0xF0, 0xE2)
WARN_LN = RGBColor(0xF0, 0xDD, 0xC0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xFA, 0xFB, 0xFC)

SANS = "Segoe UI"
MONO = "Consolas"

M = 0.55
PAGE_W, PAGE_H = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(PAGE_W)
prs.slide_height = Inches(PAGE_H)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def hline(slide, x, y, w, color=HAIR, weight=1.0):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def text(slide, x, y, w, h, runs, size=13, color=INK, font=SANS, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, caps=False,
         tracking=None, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        segs = para if isinstance(para, list) else [(para, color, bold)]
        for seg in segs:
            t, c, b = (seg + (color, bold))[:3] if isinstance(seg, tuple) else (seg, color, bold)
            r = p.add_run()
            r.text = t.upper() if caps else t
            r.font.size = Pt(size)
            r.font.name = font
            r.font.bold = b
            r.font.color.rgb = c
            if tracking is not None:
                r.font._rPr.set("spc", str(tracking))
    return tb


def badge(slide, x, y, w, label, fg=OK, bg=OK_BG, ln=OK_LN, size=12, bold=True):
    b = rect(slide, x, y, w, 0.44, fill=bg, line=ln, line_w=1.0,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = b.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = SANS
    r.font.color.rgb = fg
    return b


def new_slide(kicker, title, sub=None, page=None, total=5):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    rect(slide, 0, 0, PAGE_W / 2, 0.07, fill=PY)
    rect(slide, PAGE_W / 2, 0, PAGE_W / 2, 0.07, fill=LV)
    text(slide, M, 0.30, 10.5, 0.26, kicker, size=10.5, color=LV, font=MONO, tracking=140)
    if page is not None:
        text(slide, PAGE_W - M - 1.2, 0.30, 1.2, 0.26, f"{page} / {total}",
             size=10.5, color=MUTED, font=MONO, align=PP_ALIGN.RIGHT)
    text(slide, M, 0.56, 12.2, 0.9, title, size=28, color=INK, bold=True)
    if sub is not None:
        text(slide, M, 1.50, 12.2, 0.70, sub, size=13.5, color=MUTED, spacing=1.12)
    return slide


KICKER = "MONARCH · CONTROLS → PYTHON MIGRATION · WHAT’S NEXT · 2026-07-07"

# ======================================================================
# SLIDE 1 — the road ahead (authority ladder + the three numbers)
# ======================================================================
s = new_slide(
    KICKER,
    "The road ahead: from watching to commanding",
    [[("The shadow phase proved Python ", MUTED, False), ("decides", INK, True),
      (" correctly. The next phases give it hands — carefully: each increase in authority "
       "is built, then deliberately broken on the bench, and only granted after it fails "
       "safely and the team signs off.", MUTED, False)]],
    page=1,
)

# --- authority ladder ---
LY = 2.62
ladder = [
    ("PHASE B", "Command path", "wire the LabVIEW write\npath, then drill it", WARN, WARN_BG, WARN_LN, "UP NEXT"),
    ("PHASE C", "Bench command", "Python drives modes &\nsetpoints — no engine", PY, PY_BG, PY_LN, None),
    ("PHASE D", "Procedures", "automated sequences from\nthe team’s written steps", PY, PY_BG, PY_LN, None),
    ("PHASE E", "Commissioning", "sequences drive cold flow\n→ motoring → first fire", LV, LV_BG, LV_LN, None),
]
cw, gap = 2.62, 0.62
total_w = 4 * cw + 3 * gap
lx0 = (PAGE_W - total_w) / 2
for i, (tag, name, sub, fg, bg, ln, flag) in enumerate(ladder):
    cx = lx0 + i * (cw + gap)
    card = rect(s, cx, LY, cw, 1.42, fill=bg, line=ln, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rows = [(tag, 9, False, MONO, fg), (name, 15, True, SANS, fg)] + \
           [(line, 9.5, False, SANS, INK) for line in sub.split("\n")]
    for j, (t, sz, b, fn, col) in enumerate(rows):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(1)
        r = p.add_run(); r.text = t
        r.font.size = Pt(sz); r.font.bold = b; r.font.name = fn; r.font.color.rgb = col
    if flag:
        fb = rect(s, cx + cw / 2 - 0.55, LY - 0.19, 1.1, 0.32, fill=WARN, line=None,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        tf = fb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = flag
        r.font.size = Pt(9); r.font.bold = True; r.font.name = MONO; r.font.color.rgb = WHITE
    if i < 3:
        gx = cx + cw
        text(s, gx, LY + 0.42, gap, 0.3, "→", size=16, color=MUTED, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, gx - 0.30, LY + 0.78, gap + 0.60, 0.4, "GATE: drills +\nsign-off",
             size=7.5, color=MUTED, font=MONO, align=PP_ALIGN.CENTER, spacing=1.05)

# --- the three numbers ---
NY = 4.72
text(s, M, NY - 0.36, 8.0, 0.24, "THE NEAR-TERM WORK, IN THREE NUMBERS", size=10.5,
     color=MUTED, font=MONO, tracking=110)
tiles = [
    ("1", "LABVIEW BUILD LEFT", "the gateway write path (B3) — the only build between "
     "shadow mode and Python’s first commands; specified click-by-click", WARN),
    ("7", "FAILURE DRILLS", "we break the link every way we can think of, live — every "
     "failure must end in a rejection or a safe hold", PY),
    ("0", "AUTHORITY UNTIL THEN", "Python’s writes stay ignored (source switch on UI) "
     "until every drill passes and the team signs off", LV),
]
tw = (PAGE_W - 2 * M - 2 * 0.3) / 3
for i, (big, cap, sub, fg) in enumerate(tiles):
    tx = M + i * (tw + 0.3)
    rect(s, tx, NY, tw, 1.85, fill=PANEL, line=HAIR, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    rect(s, tx, NY, 0.07, 1.85, fill=fg)
    text(s, tx + 0.28, NY + 0.12, tw - 0.5, 0.5, big, size=27, color=fg, bold=True, wrap=False)
    text(s, tx + 0.28, NY + 0.62, tw - 0.5, 0.26, cap, size=10.5, color=MUTED, font=MONO,
         tracking=90)
    text(s, tx + 0.28, NY + 0.95, tw - 0.5, 0.85, sub, size=10.5, color=INK, spacing=1.12)

text(s, M, 6.90, 12.2, 0.4,
     [[("Head start:  ", INK, True),
       ("the Python software for every phase — commander, operator tools, sequencing "
        "engine, warning rules — is already written and tested. The near-term work is "
        "LabVIEW wiring, bench drills, and team decisions.", MUTED, False)]],
     size=11.5, spacing=1.15)

# ======================================================================
# SLIDE 2 — B3: the one LabVIEW build
# ======================================================================
s = new_slide(
    KICKER,
    "One LabVIEW build between shadow and command",
    [[("Today the gateway only ", MUTED, False), ("reads", INK, True),
      (" — it streams telemetry and politely acknowledges commands without acting. "
       "B3 teaches it to accept Python’s requests, with LabVIEW still validating every "
       "one, and a physical operator switch choosing who is in command.", MUTED, False)]],
    page=2,
)

LX, LW = M, 5.85
text(s, LX, 2.42, LW, 0.25, "WHAT GETS BUILT (3 SMALL PIECES)", size=10.5, color=MUTED,
     font=MONO, tracking=110)
hline(s, LX, 2.74, LW)
bullets = [
    [("Two new shared variables", INK, True),
     (" — a place for the operator’s inputs while Python is in command, and a "
      "UI / PYTHON source switch that ", INK, False),
     ("defaults to UI", INK, True), (".", INK, False)],
    [("The gateway command branch", INK, True),
     (" — parse each request, run a ", INK, False),
     ("6-check validation ladder", INK, True),
     (" (right command? within rate? right source? parses? in range? not clearing "
      "an e-stop?), then write — or reject with a stated reason.", INK, False)],
    [("The HMI switch + single-writer rule", INK, True),
     (" — exactly one writer of the control settings at any moment; the operator "
      "owns the switch, and only the operator.", INK, False)],
]
tb = text(s, LX, 2.92, LW, 2.6, [[("▪  ", LV, True)] + b for b in bullets],
          size=12.5, spacing=1.08)
for p in tb.text_frame.paragraphs:
    p.space_after = Pt(9)

text(s, LX, 5.85, LW, 0.25, "READY TO START", size=10.5, color=MUTED, font=MONO, tracking=110)
hline(s, LX, 6.17, LW)
tb = text(s, LX, 6.31, LW, 1.0,
          [[("▪  ", OK, True), ("Interface spec frozen with the team (ICD v0.2).", INK, False)],
           [("▪  ", OK, True),
            ("Python + simulator sides built; 13 failure tests green.", INK, False)],
           [("▪  ", OK, True),
            ("Every LabVIEW step written up node-by-node, ~half a day.", INK, False)]],
          size=11.5, spacing=1.1)
for p in tb.text_frame.paragraphs:
    p.space_after = Pt(4)

RX, RW = 6.85, 5.93
text(s, RX, 2.42, RW, 0.25, "WHO WRITES, IN EACH MODE", size=10.5, color=MUTED,
     font=MONO, tracking=110)
hline(s, RX, 2.74, RW)

def writer_panel(y, tag, tag_fg, tag_bg, tag_ln, rows):
    rect(s, RX, y, RW - 0.4, 1.28, fill=PANEL, line=HAIR, line_w=0.75,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    hdr = rect(s, RX + 0.16, y + 0.14, 2.35, 0.32, fill=tag_bg, line=tag_ln, line_w=1.0,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = hdr.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag
    r.font.size = Pt(9.5); r.font.bold = True; r.font.name = MONO; r.font.color.rgb = tag_fg
    text(s, RX + 0.22, y + 0.54, RW - 0.8, 0.7, rows, size=11, color=INK, spacing=1.18)

writer_panel(2.90, "SOURCE = UI  (DEFAULT)", LV, LV_BG, LV_LN,
             [[("Operator UI ", LV, True), ("→ writes the control settings, exactly as today", INK, False)],
              [("Python ", PY, True), ("→ every command rejected: “source is UI”", INK, False)]])
writer_panel(4.32, "SOURCE = PYTHON", PY, PY_BG, PY_LN,
             [[("Python ", PY, True), ("→ validated commands write the control settings", INK, False)],
              [("Operator UI ", LV, True), ("→ inputs become requests that Python reads", INK, False)]])

text(s, RX, 5.80, RW - 0.3, 0.8,
     [[("Either way:  ", INK, True),
       ("the 9056 state machine consumes the same variable and keeps clamping "
        "everything — the LabVIEW limiter never leaves the loop. And e-stop works "
        "from every surface, in every mode.", MUTED, False)]],
     size=11, spacing=1.15)

# ======================================================================
# SLIDE 3 — B4: the seven drills
# ======================================================================
s = new_slide(
    KICKER,
    "Nothing gets authority until it fails safely: 7 drills",
    [[("Before Python may command anything real, we try to break the link every way we "
       "can think of — ", MUTED, False), ("live on the bench", INK, True),
      (". Every failure must end the same two ways: a rejected request, or a safe hold. "
       "No exceptions, repeatably, witnessed.", MUTED, False)]],
    page=3,
)

drills = [
    ("1", "Kill Python mid-command", "cRIO clamps to SAFE within 5 s — the watchdog "
     "we just wired, re-proven through the full command path"),
    ("2", "Pull the network cable", "same safe hold; the link re-establishes itself "
     "cleanly when the cable returns"),
    ("3", "Freeze the process", "heartbeat stops changing while the socket stays up — "
     "watchdog still trips → SAFE"),
    ("4", "Send garbage", "malformed JSON is rejected with a reason; the connection "
     "survives and keeps serving telemetry"),
    ("5", "Flood commands", "more than 5 per second → rate limit trips, excess "
     "rejected — no queue, no lag"),
    ("6", "Flip the source mid-run", "UI ↔ Python handover is bumpless — values match "
     "before the switch, nothing jumps"),
    ("7", "E-stop vs Python", "e-stop from any surface always wins and latches; "
     "Python can request it but can never clear it"),
]
gy0, gh, gvs = 2.42, 0.88, 0.98
col_w = (PAGE_W - 2 * M - 0.3) / 2
for i, (num, t, sub) in enumerate(drills):
    col, row = i % 2, i // 2
    dx = M + col * (col_w + 0.3)
    dy = gy0 + row * gvs
    rect(s, dx, dy, col_w, gh, fill=PANEL, line=HAIR, line_w=0.75,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    circ = rect(s, dx + 0.14, dy + 0.14, 0.38, 0.38, fill=WARN_BG, line=WARN_LN,
                line_w=1.0, shape=MSO_SHAPE.OVAL)
    tf = circ.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.size = Pt(12); r.font.bold = True; r.font.name = MONO; r.font.color.rgb = WARN
    text(s, dx + 0.66, dy + 0.10, col_w - 0.85, 0.28, t, size=12, color=INK, bold=True)
    text(s, dx + 0.66, dy + 0.40, col_w - 0.85, 0.45, sub, size=9.5, color=MUTED, spacing=1.08)

# the 8th cell: what passing means
dx = M + 1 * (col_w + 0.3)
dy = gy0 + 3 * gvs
rect(s, dx, dy, col_w, gh, fill=OK_BG, line=OK_LN, line_w=1.0,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
text(s, dx + 0.2, dy + 0.10, col_w - 0.4, 0.28,
     [[("✔  PASS = ", OK, True), ("all seven, repeatably, live", INK, True)]], size=12)
text(s, dx + 0.2, dy + 0.40, col_w - 0.4, 0.45,
     "command → acknowledged → effect confirmed in telemetry; then joint sign-off "
     "closes Phase B", size=9.5, color=MUTED, spacing=1.08)

text(s, M, gy0 + 4 * gvs + 0.06, PAGE_W - 2 * M, 0.6,
     [[("Why so much ceremony:  ", INK, True),
       ("this gate is the safety case. It demonstrates — not argues — that a dead, hung, "
        "or misbehaving Python can only ever cause a rejected request or a safe hold.",
        MUTED, False)]],
     size=11.5, spacing=1.15)

# ======================================================================
# SLIDE 4 — Phase C: Python commands the bench
# ======================================================================
s = new_slide(
    KICKER,
    "Then Python takes the bench — supervised",
    [[("Phase C makes Python the ", MUTED, False), ("acting supervisor", INK, True),
      (" for bench operation — real authority against the real controllers, no engine "
       "hardware — while the LabVIEW clamp stays permanently underneath and the operator "
       "keeps an instant way back.", MUTED, False)]],
    page=4,
)

LX, LW = M, 5.85
text(s, LX, 2.42, LW, 0.25, "WHAT CHANGES", size=10.5, color=MUTED, font=MONO, tracking=110)
hline(s, LX, 2.74, LW)
bullets = [
    [("Python drives ", INK, False), ("mode requests and setpoints", INK, True),
     (" once per second through the command path Phase B just proved.", INK, False)],
    [("The operator keeps the LabVIEW HMI: ", INK, False),
     ("monitoring, e-stop, and the source switch", INK, True),
     (" — one flip hands control back to the UI, bumplessly.", INK, False)],
    [("The shadow comparison now runs ", INK, False), ("in reverse", INK, True),
     (": Python predicts what the LabVIEW limiter will allow, and any disagreement "
      "raises an alarm immediately.", INK, False)],
]
tb = text(s, LX, 2.92, LW, 2.4, [[("▪  ", LV, True)] + b for b in bullets],
          size=12.5, spacing=1.08)
for p in tb.text_frame.paragraphs:
    p.space_after = Pt(9)

RX, RW = 6.85, 5.93
text(s, RX, 2.42, RW, 0.25, "HOW WE’LL KNOW IT WORKS", size=10.5, color=MUTED,
     font=MONO, tracking=110)
hline(s, RX, 2.74, RW)
bullets = [
    [("Soak runs", INK, True),
     (": hours-long bench sessions with faults injected on purpose — no drift, "
      "no resource leaks, clean reconnects every time.", INK, False)],
    [("Exit test", INK, True),
     (": one scripted session — mode walks, setpoint changes, forced warnings, "
      "e-stop, kill-Python-and-recover — runs ", INK, False),
     ("entirely under Python command", INK, True),
     (" with zero unsafe outcomes and zero unexplained limiter divergences.", INK, False)],
]
tb = text(s, RX, 2.92, RW, 2.0, [[("▪  ", LV, True)] + b for b in bullets],
          size=12.5, spacing=1.08)
for p in tb.text_frame.paragraphs:
    p.space_after = Pt(9)

hline(s, M, 5.35, PAGE_W - 2 * M)
text(s, M, 5.52, 6.2, 0.25, "WHY THIS PHASE MATTERS", size=10.5, color=MUTED,
     font=MONO, tracking=110)
text(s, M, 5.84, PAGE_W - 2 * M, 1.2,
     [[("It converts trust from theory to track record. ", INK, True),
       ("Shadow mode proved the decisions; Phase B proves the failure modes; Phase C "
        "accumulates supervised operating hours. By the time automated sequences arrive "
        "(Phase D), Python isn’t a new component anymore — it’s the supervisor the bench "
        "has already been running under.", MUTED, False)]],
     size=12.5, spacing=1.2)

# ======================================================================
# SLIDE 5 — team inputs + the destination
# ======================================================================
s = new_slide(
    KICKER,
    "What we need from the team — and where this lands",
    None,
    page=5,
)

LX, LW = M, 6.6
text(s, LX, 1.65, LW, 0.25, "INPUTS THAT UNBLOCK THE NEXT PHASES", size=10.5, color=MUTED,
     font=MONO, tracking=110)
hline(s, LX, 1.97, LW)
asks = [
    ("Operating-procedure sheets (blocks Phase D)",
     "cold start, purge, motoring → light-off, shutdown, vent & recovery, misfire "
     "recovery — as informal notes or the template sheets; they become tested, "
     "automated sequences. This is the critical path."),
    ("Real numbers (blocks Phase E)",
     "warning thresholds and time windows (“X low for Y s in state Z ⇒ act”) and "
     "operating-point tables — the engines are built and waiting for values."),
    ("Two joint decisions",
     "where the UI / PYTHON switch lives on the HMI, and the sign-off review at "
     "each authority gate (B exit is the first)."),
]
ay = 2.15
for t, sub in asks:
    text(s, LX, ay, LW, 0.3, [[("▪  ", LV, True), (t, INK, True)]], size=13)
    text(s, LX + 0.25, ay + 0.32, LW - 0.25, 0.85, sub, size=11, color=MUTED, spacing=1.12)
    ay += 1.32

RX, RW = 7.65, 5.13
text(s, RX, 1.65, RW, 0.25, "THE DESTINATION", size=10.5, color=MUTED, font=MONO, tracking=110)
hline(s, RX, 1.97, RW)
dest = [
    ("Sequences drive commissioning itself", PY, PY_BG, PY_LN,
     "cold flow → motoring → first fire, each step a written, tested, repeatable "
     "procedure — not a one-off manual run"),
    ("Non-combustion sequences go first", LV, LV_BG, LV_LN,
     "venting, purge, working-fluid checks, thermal warm-up — already drafted, and "
     "exactly what commissioning needs earliest"),
    ("Every capability earns its place", OK, OK_BG, OK_LN,
     "nothing is unlocked until its bench drill passes — the same rule from day one, "
     "all the way up"),
]
dy = 2.15
for t, fg, bg, ln, sub in dest:
    rect(s, RX, dy, RW - 0.1, 1.18, fill=bg, line=ln, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    text(s, RX + 0.2, dy + 0.12, RW - 0.5, 0.28, t, size=12, color=fg, bold=True)
    text(s, RX + 0.2, dy + 0.44, RW - 0.5, 0.65, sub, size=10, color=INK, spacing=1.1)
    dy += 1.32

# safety footer
fy = 6.35
fb = rect(s, M, fy, PAGE_W - 2 * M, 0.72, fill=LV_BG, line=LV_LN, line_w=1.0,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
tf = fb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.word_wrap = True
tf.margin_left = Inches(0.22)
tf.margin_right = Inches(0.22)
p = tf.paragraphs[0]
for seg, c, b in [("The rule that never changes:  ", LV, True),
                  ("LabVIEW/cRIO keeps the hardware, the interlocks and the safe fallback — "
                   "permanently, through every phase on this roadmap. Python only ever asks; "
                   "LabVIEW checks every request; if Python disappears, the system safe-holds.",
                   INK, False)]:
    r = p.add_run(); r.text = seg
    r.font.size = Pt(11.5); r.font.bold = b; r.font.name = SANS; r.font.color.rgb = c

out = HERE / "migration-next-steps.pptx"
prs.save(str(out))
print("wrote", out)
