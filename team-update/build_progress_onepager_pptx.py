"""Build the July-7 progress update as a SINGLE 16:9 slide (editable PowerPoint).

One-slide version of build_progress_pptx.py, laid out like the original
hello-VI one-slider: header, plain-language stories on the left with the
shadow-mode concept diagram, proof numbers on the right, and a phase-status
strip along the bottom.

    python team-update/build_progress_onepager_pptx.py
        ->  team-update/migration-progress-onepager.pptx
"""

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = pathlib.Path(__file__).parent

# palette (same as the other team-update decks)
INK = RGBColor(0x16, 0x22, 0x2D)
MUTED = RGBColor(0x5C, 0x6B, 0x78)
HAIR = RGBColor(0xE2, 0xE7, 0xEC)
PY = RGBColor(0x5B, 0x46, 0xC9)
PY_BG = RGBColor(0xED, 0xEA, 0xFA)
PY_LN = RGBColor(0xD6, 0xCE, 0xF4)
LV = RGBColor(0x0E, 0x76, 0x7E)
LV_BG = RGBColor(0xE0, 0xF0, 0xF0)
LV_LN = RGBColor(0xC4, 0xE2, 0xE2)
OK = RGBColor(0x1C, 0x91, 0x4F)
OK_BG = RGBColor(0xE3, 0xF4, 0xEB)
OK_LN = RGBColor(0xBF, 0xE6, 0xCE)
WARN = RGBColor(0xB4, 0x5B, 0x0E)
WARN_BG = RGBColor(0xFB, 0xF0, 0xE2)
WARN_LN = RGBColor(0xF0, 0xDD, 0xC0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xFA, 0xFB, 0xFC)

SANS = "Segoe UI"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


def rect(x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def hline(x, y, w, color=HAIR, weight=1.0):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def text(x, y, w, h, runs, size=13, color=INK, font=SANS, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.0, caps=False, tracking=None, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        segs = para if isinstance(para, list) else [(para, color, bold)]
        for seg in segs:
            t, c, b = (seg + (color, bold))[:3] if isinstance(seg, tuple) else (seg, color, bold)
            r = p.add_run()
            r.text = t.upper() if caps else t
            r.font.size = Pt(size)
            r.font.name = font
            r.font.bold = b
            r.font.color.rgb = c
            if tracking is not None:
                r.font._rPr.set("spc", str(tracking))
    return tb


def node(x, y, w, h, tag, name, sub, bg, ln, fg):
    sp = rect(x, y, w, h, fill=bg, line=ln, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rows = [(tag.upper(), 9, False, MONO), (name, 14, True, SANS), (sub, 9, False, SANS)]
    for i, (t, s, b, fn) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = t
        r.font.size = Pt(s)
        r.font.bold = b
        r.font.name = fn
        r.font.color.rgb = fg
    return sp


M = 0.55

# --- accent bar ---
rect(0, 0, 6.6665, 0.07, fill=PY)
rect(6.6665, 0, 6.6665, 0.07, fill=LV)

# --- header ---
text(M, 0.30, 10.5, 0.26, "MONARCH · CONTROLS → PYTHON MIGRATION · UPDATE 2 · 2026-07-07",
     size=10.5, color=LV, font=MONO, tracking=140)
text(M, 0.56, 12.2, 0.9, "From “hello” to a proven brain — five days of progress",
     size=26, color=INK, bold=True)
text(M, 1.42, 12.2, 0.70,
     [[("Since the July 2 connectivity test: the first supervisory state machine now runs in "
        "Python and ", MUTED, False),
       ("matches the live LabVIEW system decision-for-decision", INK, True),
       (" — and along the way we found and fixed a real safety gap that predates this project.",
        MUTED, False)]],
     size=13.5, color=MUTED, spacing=1.12)

# ============ LEFT COLUMN — the three stories ============
LX, LW = M, 5.85
text(LX, 2.18, LW, 0.25, "WHAT HAPPENED, IN PLAIN TERMS", size=10.5, color=MUTED,
     font=MONO, tracking=110)
hline(LX, 2.50, LW)

bullets = [
    [("Brain ported & proven.  ", INK, True),
     ("The 9056’s state machine (SAFE→FIRING) and warning policy now run in Python, "
      "validated in shadow mode against the live system — every state, e-stop, "
      "overrides, warning clamp.", INK, False)],
    [("Safety gap found & fixed.  ", INK, True),
     ("The cRIO detected a lost PC but the response was never wired (original code). "
      "Now it clamps to SAFE in 5 s — proven live, pull-the-plug.", INK, False)],
    [("Built ahead.  ", INK, True),
     ("The Python side of every remaining phase — command path, operator tools, "
      "sequencing, warning rules — is written, simulated, and tested.", INK, False)],
]
tb = text(LX, 2.66, LW, 2.2, [[("▪  ", LV, True)] + b for b in bullets], size=11.5, spacing=1.05)
for p in tb.text_frame.paragraphs:
    p.space_after = Pt(6)

# shadow-mode concept diagram (one horizontal row)
DY = 5.10
text(LX, DY - 0.26, LW, 0.24, "SHADOW MODE · SAME DATA, TWO BRAINS, EVERY FRAME COMPARED",
     size=8.5, color=MUTED, font=MONO, align=PP_ALIGN.CENTER, tracking=60)
node(LX, DY, 1.70, 0.88, "In charge", "LabVIEW", "decides & acts", LV_BG, LV_LN, LV)
text(LX + 1.70, DY, 0.36, 0.88, "→", size=15, color=MUTED, bold=True,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
node(LX + 2.06, DY, 1.70, 0.88, "Shadow", "Python", "sends nothing", PY_BG, PY_LN, PY)
text(LX + 3.76, DY, 0.36, 0.88, "→", size=15, color=MUTED, bold=True,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
node(LX + 4.12, DY, 1.70, 0.88, "Compared", "✔ agree", "every frame", OK_BG, OK_LN, OK)

# ============ RIGHT COLUMN — the proof ============
RX, RW = 6.85, 5.93
text(RX, 2.18, RW, 0.25, "THE PROOF", size=10.5, color=MUTED, font=MONO, tracking=110)
hline(RX, 2.50, RW)

stats = [
    ("100%", PY, "live agreement, all 5 states",
     "Python matched LabVIEW across every input we could throw at it (one differing frame "
     "traced to a telemetry timing artifact, not the logic)"),
    ("5 s", OK, "PC lost → cRIO safe-holds",
     "live drop episode: heartbeat stops → SAFE hold → controlled step-by-step recovery; "
     "Python’s shadow agreed on every frame"),
    ("144+", LV, "automated tests, on every change",
     "all ported logic pinned by a test suite run in CI, so it can’t silently drift"),
]
sy = 2.66
for big, fg, cap, sub in stats:
    rect(RX, sy, RW, 1.02, fill=PANEL, line=HAIR, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    rect(RX, sy, 0.06, 1.02, fill=fg)
    text(RX + 0.22, sy + 0.10, 1.30, 0.5, big, size=22, color=fg, bold=True, wrap=False)
    text(RX + 1.55, sy + 0.13, RW - 1.75, 0.28, cap, size=11.5, color=INK, bold=True)
    text(RX + 1.55, sy + 0.42, RW - 1.75, 0.55, sub, size=9.5, color=MUTED, spacing=1.08)
    sy += 1.14

# ============ BOTTOM BAND — phase strip + next ============
BY = 6.24
hline(M, BY - 0.06, 13.333 - 2 * M)
text(M, BY + 0.02, 2.2, 0.22, "PHASES A–E", size=9, color=MUTED, font=MONO, tracking=110)

phases = [
    ("A  shadow brain — validated live", OK, OK_BG, OK_LN),
    ("B  command path — LabVIEW next", WARN, WARN_BG, WARN_LN),
    ("C  operator tools — built", PY, PY_BG, PY_LN),
    ("D  sequencing — built", PY, PY_BG, PY_LN),
    ("E  warning rules — built", PY, PY_BG, PY_LN),
]
cx = M
for label, fg, bg, ln in phases:
    cw = 0.26 + 0.082 * len(label)
    ch = rect(cx, BY + 0.24, cw, 0.34, fill=bg, line=ln, line_w=1.0,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = ch.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(9.5); r.font.name = MONO; r.font.bold = True; r.font.color.rgb = fg
    cx += cw + 0.13

text(11.9, BY + 0.02, 0.9, 0.3, "2026-07-07", size=9, color=MUTED, font=MONO,
     align=PP_ALIGN.RIGHT)
text(M, BY + 0.70, 12.2, 0.55,
     [[("Why it matters:  ", INK, True),
       ("the riskiest step — proving Python can faithfully reproduce the control logic against "
        "the real system — is done, and “Python-offline = safe hold” is now enforced, not assumed.   ",
        MUTED, False),
       ("Next:  ", INK, True),
       ("wire the LabVIEW command write path (B3), then the bench failure drills (B4).",
        MUTED, False)]],
     size=11, spacing=1.15)

out = HERE / "migration-progress-onepager.pptx"
prs.save(str(out))
print("wrote", out)
