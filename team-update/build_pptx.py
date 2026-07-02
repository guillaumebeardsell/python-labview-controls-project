"""Build the one-slide team update as an editable 16:9 PowerPoint.

Mirrors team-update/hello-vi-slide.html: header, plain-language points + a
Python<->LabVIEW concept diagram, the console PASS screenshot as proof, an
"under the hood" block-diagram thumbnail, and a result strip.

    python team-update/build_pptx.py   ->   team-update/hello-vi-update.pptx
"""

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

HERE = pathlib.Path(__file__).parent

# palette (from the HTML)
INK = RGBColor(0x16, 0x22, 0x2D)
MUTED = RGBColor(0x5C, 0x6B, 0x78)
HAIR = RGBColor(0xE2, 0xE7, 0xEC)
PY = RGBColor(0x5B, 0x46, 0xC9)
PY_BG = RGBColor(0xED, 0xEA, 0xFA)
PY_LN = RGBColor(0xD6, 0xCE, 0xF4)
LV = RGBColor(0x0E, 0x76, 0x7E)
LV_BG = RGBColor(0xE0, 0xF0, 0xF0)
LV_LN = RGBColor(0xC4, 0xE2, 0xE2)
OK = RGBColor(0x1C, 0x91, 0x4F)
OK_BG = RGBColor(0xE3, 0xF4, 0xEB)
OK_LN = RGBColor(0xBF, 0xE6, 0xCE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xFA, 0xFB, 0xFC)

SANS = "Segoe UI"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


def rect(x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def hline(x, y, w, color=HAIR, weight=1.0):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def text(x, y, w, h, runs, size=13, color=INK, font=SANS, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.0, caps=False, tracking=None, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        segs = para if isinstance(para, list) else [(para, color, bold)]
        for seg in segs:
            t, c, b = (seg + (color, bold))[:3] if isinstance(seg, tuple) else (seg, color, bold)
            r = p.add_run()
            r.text = t.upper() if caps else t
            r.font.size = Pt(size)
            r.font.name = font
            r.font.bold = b
            r.font.color.rgb = c
            if tracking is not None:
                r.font._rPr.set("spc", str(tracking))
    return tb


def node(x, y, w, h, tag, name, sub, bg, ln, fg):
    sp = rect(x, y, w, h, fill=bg, line=ln, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rows = [(tag.upper(), 9, False, MONO), (name, 15, True, SANS), (sub, 9.5, False, SANS)]
    for i, (t, s, b, fn) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = t
        r.font.size = Pt(s)
        r.font.bold = b
        r.font.name = fn
        r.font.color.rgb = fg
    return sp


M = 0.55

# --- accent bar ---
rect(0, 0, 6.6665, 0.07, fill=PY)
rect(6.6665, 0, 6.6665, 0.07, fill=LV)

# --- header ---
text(M, 0.30, 9.0, 0.26, "MONARCH · CONTROLS → PYTHON MIGRATION",
     size=10.5, color=LV, font=MONO, tracking=140)
text(M, 0.56, 12.2, 0.9, "The new Python layer is talking to LabVIEW — live",
     size=29, color=INK, bold=True)
text(M, 1.52, 8.9, 0.70,
     [[("First end-to-end connection test ", MUTED, False), ("passed", INK, True),
       (": the two systems exchanged data both ways, once a second, running side by side "
        "on the control-room PC — with no change to the safety-critical LabVIEW control.",
        MUTED, False)]],
     size=13.5, color=MUTED, spacing=1.12)

# ============ LEFT COLUMN ============
LX, LW = M, 5.85
text(LX, 2.32, LW, 0.25, "THE TEST, IN PLAIN TERMS", size=10.5, color=MUTED, font=MONO, tracking=110)
hline(LX, 2.64, LW)

bullets = [
    [("Built a minimal ", INK, False), ("“hello” link", INK, True),
     (" between the new Python program and the existing LabVIEW system.", INK, False)],
    [("LabVIEW streams a live value to Python ", INK, False), ("once per second", INK, True),
     ("; Python sends commands back and LabVIEW ", INK, False),
     ("acknowledges every one", INK, True), (".", INK, False)],
    [("Both ran on the same PC at once — and ", INK, False),
     ("either side can restart", INK, True), (" and the link re-establishes itself.", INK, False)],
]
tb = text(LX, 2.8, LW, 1.75, [[("▪  ", LV, True)] + b for b in bullets],
          size=13, spacing=1.06)
for p in tb.text_frame.paragraphs:
    p.space_after = Pt(9)

# concept diagram
text(LX, 4.58, LW, 0.24, "ONE CONTROL-ROOM PC · BOTH RUNNING",
     size=9, color=MUTED, font=MONO, align=PP_ALIGN.CENTER, tracking=90)
node(LX, 4.88, 2.0, 1.02, "New", "Python", "supervisory layer", PY_BG, PY_LN, PY)
node(LX + 3.85, 4.88, 2.0, 1.02, "Existing", "LabVIEW", "engine control", LV_BG, LV_LN, LV)
text(LX + 2.05, 4.9, 1.75, 1.0,
     [[("live data →", LV, True)], [("← commands", PY, True)], [("♥ heartbeat", MUTED, False)]],
     size=9.5, font=MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.35)

# ============ RIGHT COLUMN ============
RX, RW = 6.85, 5.93
text(RX, 2.32, RW, 0.25, "PROOF IT WORKED", size=10.5, color=MUTED, font=MONO, tracking=110)
hline(RX, 2.64, RW)

badge = rect(RX, 2.78, 4.75, 0.44, fill=OK_BG, line=OK_LN, line_w=1.0,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
btf = badge.text_frame
btf.vertical_anchor = MSO_ANCHOR.MIDDLE
bp = btf.paragraphs[0]
bp.alignment = PP_ALIGN.CENTER
br = bp.add_run()
br.text = "✔  PASS — two-way communication confirmed"
br.font.size = Pt(12)
br.font.bold = True
br.font.name = SANS
br.font.color.rgb = OK

# CMD screenshot (ratio 1.733)
cmd_w = 4.42
cmd_h = cmd_w / 1.733
rect(RX, 3.30, cmd_w, cmd_h, fill=RGBColor(0x0C, 0x0C, 0x0C), line=HAIR)
slide.shapes.add_picture(str(HERE / "CMD-Output.png"), Inches(RX), Inches(3.30),
                         Inches(cmd_w), Inches(cmd_h))
text(RX, 3.30 + cmd_h + 0.05, RW, 0.3,
     "Python console: live values every second → RESULT: PASS",
     size=9.5, color=MUTED, font=MONO)
# small "LabVIEW side connected" proof, right of the console
fp_x = RX + cmd_w + 0.14
fp_w, fp_h = 1.30, 1.30 / 3.3
slide.shapes.add_picture(str(HERE / "Hello-VI_Front-Panel.png"),
                         Inches(fp_x), Inches(3.30), Inches(fp_w), Inches(fp_h))
text(fp_x, 3.30 + fp_h + 0.05, 1.30, 0.7,
     "LabVIEW side: Client Connected", size=8.5, color=MUTED, font=MONO, spacing=1.1)

# ============ BOTTOM BAND ============
BY = 6.30
hline(M, BY - 0.06, 13.333 - 2 * M)

# under the hood (left): block-diagram thumbnail
text(M, BY, 2.5, 0.22, "UNDER THE HOOD", size=9, color=MUTED, font=MONO, tracking=110)
bd_w = 2.0
bd_h = bd_w / 2.54
rect(M, BY + 0.25, bd_w, bd_h, fill=WHITE, line=HAIR)
slide.shapes.add_picture(str(HERE / "Hello-VI_Block-Diagram.png"),
                         Inches(M + 0.05), Inches(BY + 0.30), Inches(bd_w - 0.10), Inches(bd_h - 0.10))
text(M + bd_w + 0.16, BY + 0.42, 1.15, 0.6,
     "the LabVIEW\ntest program", size=9, color=MUTED, font=MONO, spacing=1.1)

# result chips + what's next (right)
chips = ["✓ Two-way link", "✓ ~1 Hz updates", "✓ Auto-reconnect"]
cx = 4.55
for c in chips:
    cw = 0.20 + 0.093 * len(c)
    ch = rect(cx, BY + 0.02, cw, 0.34, fill=WHITE, line=HAIR, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    ctf = ch.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    r0 = cp.add_run(); r0.text = c[0] + " "; r0.font.color.rgb = OK
    r1 = cp.add_run(); r1.text = c[2:]; r1.font.color.rgb = INK
    for r in (r0, r1):
        r.font.size = Pt(11); r.font.name = MONO
    cx += cw + 0.14

text(11.5, BY + 0.03, 1.28, 0.3, "2026-07-02", size=9, color=MUTED, font=MONO, align=PP_ALIGN.RIGHT)
text(4.55, BY + 0.52, 8.05, 0.55,
     [[("Why it matters:  ", INK, True),
       ("connecting the two systems was the biggest unknown — now proven.", MUTED, False)],
      [("Next:  ", INK, True), ("move the first supervisory state machine into Python.", MUTED, False)]],
     size=11, color=MUTED, spacing=1.18)

out = HERE / "hello-vi-update.pptx"
prs.save(str(out))
print("wrote", out)
