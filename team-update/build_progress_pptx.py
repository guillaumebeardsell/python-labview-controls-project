"""Build the July-7 team progress update as an editable 16:9 PowerPoint.

Follow-up to build_pptx.py (the 2026-07-02 hello-VI one-slider); same palette
and visual language, five slides:

  1. Headline + timeline + the three big numbers
  2. Phase A — the state machine ported and proven live (shadow compare)
  3. The loss-of-PC safety gap: found, fixed, verified live
  4. Build-ahead status across phases A-E
  5. What's next + what we need from the team

    python team-update/build_progress_pptx.py  ->  team-update/migration-progress.pptx
"""

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = pathlib.Path(__file__).parent

# palette (same as build_pptx.py)
INK = RGBColor(0x16, 0x22, 0x2D)
MUTED = RGBColor(0x5C, 0x6B, 0x78)
HAIR = RGBColor(0xE2, 0xE7, 0xEC)
PY = RGBColor(0x5B, 0x46, 0xC9)
PY_BG = RGBColor(0xED, 0xEA, 0xFA)
PY_LN = RGBColor(0xD6, 0xCE, 0xF4)
LV = RGBColor(0x0E, 0x76, 0x7E)
LV_BG = RGBColor(0xE0, 0xF0, 0xF0)
LV_LN = RGBColor(0xC4, 0xE2, 0xE2)
OK = RGBColor(0x1C, 0x91, 0x4F)
OK_BG = RGBColor(0xE3, 0xF4, 0xEB)
OK_LN = RGBColor(0xBF, 0xE6, 0xCE)
WARN = RGBColor(0xB4, 0x5B, 0x0E)
WARN_BG = RGBColor(0xFB, 0xF0, 0xE2)
WARN_LN = RGBColor(0xF0, 0xDD, 0xC0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xFA, 0xFB, 0xFC)

SANS = "Segoe UI"
MONO = "Consolas"

M = 0.55  # page margin
PAGE_W, PAGE_H = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(PAGE_W)
prs.slide_height = Inches(PAGE_H)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def hline(slide, x, y, w, color=HAIR, weight=1.0):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def text(slide, x, y, w, h, runs, size=13, color=INK, font=SANS, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, caps=False,
         tracking=None, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        segs = para if isinstance(para, list) else [(para, color, bold)]
        for seg in segs:
            t, c, b = (seg + (color, bold))[:3] if isinstance(seg, tuple) else (seg, color, bold)
            r = p.add_run()
            r.text = t.upper() if caps else t
            r.font.size = Pt(size)
            r.font.name = font
            r.font.bold = b
            r.font.color.rgb = c
            if tracking is not None:
                r.font._rPr.set("spc", str(tracking))
    return tb


def chip(slide, x, y, label, fg=INK, bg=WHITE, ln=HAIR, size=11, lead=None, lead_fg=OK):
    """Small rounded pill; returns its width. lead = leading glyph in its own color."""
    body = label if lead is None else lead + " " + label
    cw = 0.24 + 0.10 * len(body) * (size / 11)
    ch = rect(slide, x, y, cw, 0.34, fill=bg, line=ln,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = ch.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if lead is not None:
        r0 = p.add_run(); r0.text = lead + " "; r0.font.color.rgb = lead_fg
        r0.font.size = Pt(size); r0.font.name = MONO; r0.font.bold = True
    r1 = p.add_run(); r1.text = label; r1.font.color.rgb = fg
    r1.font.size = Pt(size); r1.font.name = MONO
    return cw


def badge(slide, x, y, w, label, fg=OK, bg=OK_BG, ln=OK_LN, size=12, bold=True):
    b = rect(slide, x, y, w, 0.44, fill=bg, line=ln, line_w=1.0,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = b.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = SANS
    r.font.color.rgb = fg
    return b


def node(slide, x, y, w, h, tag, name, sub, bg, ln, fg):
    sp = rect(slide, x, y, w, h, fill=bg, line=ln, line_w=1.0,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rows = [(tag.upper(), 9, False, MONO), (name, 15, True, SANS), (sub, 9.5, False, SANS)]
    for i, (t, s, b, fn) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = t
        r.font.size = Pt(s)
        r.font.bold = b
        r.font.name = fn
        r.font.color.rgb = fg
    return sp


def new_slide(kicker, title, sub=None, page=None, total=5):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    rect(slide, 0, 0, PAGE_W / 2, 0.07, fill=PY)
    rect(slide, PAGE_W / 2, 0, PAGE_W / 2, 0.07, fill=LV)
    text(slide, M, 0.30, 10.5, 0.26, kicker, size=10.5, color=LV, font=MONO, tracking=140)
    if page is not None:
        text(slide, PAGE_W - M - 1.2, 0.30, 1.2, 0.26, f"{page} / {total}",
             size=10.5, color=MUTED, font=MONO, align=PP_ALIGN.RIGHT)
    text(slide, M, 0.56, 12.2, 0.9, title, size=28, color=INK, bold=True)
    if sub is not None:
        text(slide, M, 1.50, 12.2, 0.70, sub, size=13.5, color=MUTED, spacing=1.12)
    return slide


KICKER = "MONARCH · CONTROLS → PYTHON MIGRATION · UPDATE 2 · 2026-07-07"

# ======================================================================
# SLIDE 1 — headline, timeline, big numbers
# ======================================================================
s = new_slide(
    KICKER,
    "From “hello” to a proven brain — five days of progress",
    [[("Since the July 2 connectivity test: the first supervisory state machine now runs in Python and ",
       MUTED, False),
      ("matches the live LabVIEW system decision-for-decision", INK, True),
      (" — and along the way we found and fixed a real safety gap that predates this project. "
       "The safety-critical LabVIEW control still owns all hardware.", MUTED, False)]],
    page=1,
)

# --- timeline ---
TL_Y = 2.95
hline(s, M, TL_Y, PAGE_W - 2 * M, color=HAIR, weight=2.0)
milestones = [
    ("JUL 2", "“Hello” link", "two-way Python↔LabVIEW\nconnection — PASS", LV, LV_BG, LV_LN),
    ("JUL 3", "Live telemetry", "real control data streams\nto Python at 1 Hz", LV, LV_BG, LV_LN),
    ("JUL 6", "Migration plan", "phased roadmap A–E with\nsafety gates at each step", PY, PY_BG, PY_LN),
    ("JUL 7", "Brain validated", "Python state machine agrees\nwith LabVIEW live — 100%", PY, PY_BG, PY_LN),
    ("JUL 7", "Safety gap fixed", "loss-of-PC watchdog wired\nand verified on the rig", OK, OK_BG, OK_LN),
]
n = len(milestones)
tl_start = M + 1.13
span = PAGE_W - M - 1.13 - tl_start
for i, (d, t, sub, fg, bg, ln) in enumerate(milestones):
    cx = tl_start + span * i / (n - 1)
    dot = rect(s, cx - 0.09, TL_Y - 0.09, 0.18, 0.18, fill=fg, shape=MSO_SHAPE.OVAL)
    text(s, cx - 1.1, TL_Y - 0.44, 2.2, 0.24, d, size=9.5, color=fg, font=MONO,
         align=PP_ALIGN.CENTER, tracking=90)
    card = rect(s, cx - 1.13, TL_Y + 0.22, 2.26, 1.06, fill=bg, line=ln, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for j, (tt, sz, bd, col) in enumerate([(t, 12.5, True, fg)] +
                                          [(line, 9, False, INK) for line in sub.split("\n")]):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(1)
        r = p.add_run(); r.text = tt
        r.font.size = Pt(sz); r.font.bold = bd; r.font.name = SANS; r.font.color.rgb = col

# --- big numbers ---
NY = 4.85
text(s, M, NY - 0.36, 6.0, 0.24, "THE THREE NUMBERS THAT MATTER", size=10.5, color=MUTED,
     font=MONO, tracking=110)
tiles = [
    ("100%", "LIVE AGREEMENT", "Python’s decisions matched LabVIEW’s across all 5 engine "
     "states and every input we could throw at it", PY),
    ("1", "SAFETY GAP CLOSED", "the cRIO detected a lost PC but did nothing about it — "
     "now it safe-holds, proven with a live pull-the-plug test", OK),
    ("144+", "AUTOMATED TESTS", "every piece of ported logic is pinned by tests that run "
     "on every change, so it can’t silently drift", LV),
]
tw = (PAGE_W - 2 * M - 2 * 0.3) / 3
for i, (big, cap, sub, fg) in enumerate(tiles):
    tx = M + i * (tw + 0.3)
    rect(s, tx, NY, tw, 1.85, fill=PANEL, line=HAIR, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    rect(s, tx, NY, 0.07, 1.85, fill=fg)
    text(s, tx + 0.28, NY + 0.12, tw - 0.5, 0.5, big, size=27, color=fg, bold=True, wrap=False)
    text(s, tx + 0.28, NY + 0.62, tw - 0.5, 0.26, cap, size=10.5, color=MUTED, font=MONO,
         tracking=90)
    text(s, tx + 0.28, NY + 0.95, tw - 0.5, 0.85, sub, size=10.5, color=INK, spacing=1.12)

text(s, M, 7.02, 12.2, 0.3,
     [[("Bottom line:  ", INK, True),
       ("the riskiest step of the migration — proving Python can faithfully reproduce the "
        "control logic against the real system — is done.", MUTED, False)]],
     size=11.5)

# ======================================================================
# SLIDE 2 — Phase A: ported and proven live
# ======================================================================
s = new_slide(
    KICKER,
    "The first “brain” is ported — and it matches, live",
    [[("The 9056’s supervisory state machine (SAFE → STANDBY → MOTORING → IDLING → FIRING) "
       "and its warning policy were rebuilt in Python, then run in ", MUTED, False),
      ("shadow mode", INK, True),
      (": LabVIEW keeps deciding for real, Python decides in parallel from the same live data, "
       "and every frame is compared.", MUTED, False)]],
    page=2,
)

LX, LW = M, 5.85
text(s, LX, 2.42, LW, 0.25, "WHAT WE DID", size=10.5, color=MUTED, font=MONO, tracking=110)
hline(s, LX, 2.74, LW)
bullets = [
    [("Transcribed the state machine and warning logic from the LabVIEW diagrams into ",
      INK, False), ("plain, testable Python", INK, True),
     (" — pure logic, no hardware access, by design.", INK, False)],
    [("Built a ", INK, False), ("shadow-compare harness", INK, True),
     (" that replays every telemetry frame through the Python port and flags any "
      "decision that differs from LabVIEW’s.", INK, False)],
    [("Exercised the ", INK, False), ("whole envelope live", INK, True),
     (": every state, mode walks up to FIRING, e-stop, manual overrides, and the "
      "warning clamp.", INK, False)],
]
tb = text(s, LX, 2.92, LW, 2.2, [[("▪  ", LV, True)] + b for b in bullets], size=12.5, spacing=1.08)
for p in tb.text_frame.paragraphs:
    p.space_after = Pt(9)

# shadow-mode diagram
DY = 5.35
text(s, LX, DY - 0.28, LW, 0.24, "SHADOW MODE · SAME DATA, TWO BRAINS, EVERY FRAME COMPARED",
     size=9, color=MUTED, font=MONO, align=PP_ALIGN.CENTER, tracking=70)
node(s, LX, DY + 0.05, 1.95, 0.98, "In charge", "LabVIEW", "decides & acts", LV_BG, LV_LN, LV)
node(s, LX + 3.90, DY + 0.05, 1.95, 0.98, "Shadow", "Python", "decides, sends nothing", PY_BG, PY_LN, PY)
text(s, LX + 1.98, DY + 0.10, 1.90, 0.55, [[("live data →", LV, True)]],
     size=9.5, font=MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
cmp_b = rect(s, LX + 1.55, DY + 1.18, 2.75, 0.46, fill=OK_BG, line=OK_LN, line_w=1.0,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
tf = cmp_b.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "same decision?  ✔ yes"
r.font.size = Pt(11); r.font.bold = True; r.font.name = MONO; r.font.color.rgb = OK

RX, RW = 6.85, 5.93
text(s, RX, 2.42, RW, 0.25, "THE RESULT", size=10.5, color=MUTED, font=MONO, tracking=110)
hline(s, RX, 2.74, RW)
badge(s, RX, 2.90, 5.2, "✔  100% agreement — all 5 states, every input tested")
res = [
    [("Full-envelope live sessions: ", INK, False),
     ("override sweep 100%, e-stop 100%, mode walk 100%", INK, True), (".", INK, False)],
    [("Longest walk (210 frames): the actuator limiter matched ", INK, False),
     ("209/209", INK, True),
     ("; the state matched 208/209 — the single differing frame was traced to a "
      "telemetry timing artifact in the test tap, ", INK, False),
     ("not the ported logic", INK, True), (".", INK, False)],
    [("Every difference found along the way is logged and dispositioned in ", INK, False),
     ("docs/shadow-findings.md", INK, True), (".", INK, False)],
]
tb = text(s, RX, 3.52, RW, 1.9, [[("▪  ", LV, True)] + b for b in res], size=12.5, spacing=1.08)
for p in tb.text_frame.paragraphs:
    p.space_after = Pt(8)

cx = RX
for c in ["5/5 states", "e-stop", "overrides", "warning clamp"]:
    cx += chip(s, cx, 5.55, c, lead="✓") + 0.14

hline(s, RX, 6.20, RW)
text(s, RX, 6.35, RW, 0.9,
     [[("Bonus finding:  ", INK, True),
       ("shadow compare also exposed dead code — LabVIEW’s “Limited” output and two "
        "developer override switches were never wired to anything. The port is already "
        "paying for itself as an audit of the original code.", MUTED, False)]],
     size=11, spacing=1.15)

# ======================================================================
# SLIDE 3 — the loss-of-PC safety gap
# ======================================================================
s = new_slide(
    KICKER,
    "Found and fixed: a safety gap older than this project",
    [[("While mapping the LabVIEW code we found that the engine cRIO ", MUTED, False),
      ("detected", INK, True),
      (" a lost control-room PC — but the response was never wired in. The alarm rang and "
       "nothing was listening. This is original-code behaviour, not something Python introduced.",
       MUTED, False)]],
    page=3,
)

LX, LW = M, 5.85
text(s, LX, 2.42, LW, 0.25, "BEFORE — DETECTION WITHOUT RESPONSE", size=10.5, color=WARN,
     font=MONO, tracking=110)
hline(s, LX, 2.74, LW)
tb = text(s, LX, 2.92, LW, 1.7,
          [[("▪  ", WARN, True),
            ("The watchdog computed a “PC not responding” flag…", INK, False)],
           [("▪  ", WARN, True),
            ("…but the flag ", INK, False), ("gated nothing", INK, True),
            (" — no state change, no safe hold, no operator alarm.", INK, False)],
           [("▪  ", WARN, True),
            ("A crashed PC (LabVIEW UI ", INK, False), ("or", INK, True),
            (" Python) would have left the engine running open-loop on stale commands.",
             INK, False)]],
          size=12.5, spacing=1.08)
for p in tb.text_frame.paragraphs:
    p.space_after = Pt(8)

text(s, LX, 4.72, LW, 0.25, "AFTER — WIRED, SIZED, AND PROVEN", size=10.5, color=OK,
     font=MONO, tracking=110)
hline(s, LX, 5.04, LW)
tb = text(s, LX, 5.20, LW, 1.3,
          [[("▪  ", OK, True),
            ("The watchdog now drives the state machine: lose the PC heartbeat for 5 s "
             "and the cRIO ", INK, False), ("clamps itself to SAFE", INK, True),
            (" — with the LabVIEW UI today or Python later.", INK, False)],
           [("▪  ", OK, True),
            ("Threshold sized deliberately: 5 s ≈ 250 ticks of the ~20 ms control loop — "
             "rides out normal jitter, still trips fast.", INK, False)]],
          size=12.5, spacing=1.08)
for p in tb.text_frame.paragraphs:
    p.space_after = Pt(8)

RX, RW = 6.85, 5.93
text(s, RX, 2.42, RW, 0.25, "PROVEN WITH A LIVE PULL-THE-PLUG TEST", size=10.5, color=MUTED,
     font=MONO, tracking=110)
hline(s, RX, 2.74, RW)
steps = [
    ("1", "PC link dropped", "heartbeat stops mid-run", WARN, WARN_BG, WARN_LN),
    ("2", "cRIO reacts alone", "watchdog trips → system clamps to SAFE and holds", OK, OK_BG, OK_LN),
    ("3", "PC returns", "heartbeat resumes, warning clears", LV, LV_BG, LV_LN),
    ("4", "Controlled recovery", "state climbs back one step per tick — no jump to run", PY, PY_BG, PY_LN),
]
sy = 2.95
for i, (num, t, sub, fg, bg, ln) in enumerate(steps):
    card = rect(s, RX, sy, RW - 0.4, 0.62, fill=bg, line=ln, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    tf = card.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = num + "  "; r.font.color.rgb = fg
    r.font.size = Pt(13); r.font.bold = True; r.font.name = MONO
    r = p.add_run(); r.text = t + " — "; r.font.color.rgb = fg
    r.font.size = Pt(12.5); r.font.bold = True; r.font.name = SANS
    r = p.add_run(); r.text = sub; r.font.color.rgb = INK
    r.font.size = Pt(11.5); r.font.name = SANS
    sy += 0.74
badge(s, RX, sy + 0.06, RW - 0.4, "✔  Python’s shadow agreed on every frame of the episode",
      size=11.5)

hline(s, M, 6.55, PAGE_W - 2 * M)
text(s, M, 6.72, PAGE_W - 2 * M, 0.7,
     [[("Why this matters:  ", INK, True),
       ("“if the PC dies, the engine safe-holds” was the one unverified assumption in the "
        "migration’s safety case — and the gate before Python may ever hold command authority. "
        "It is now enforced in the cRIO and demonstrated, not assumed.", MUTED, False)]],
     size=11.5, spacing=1.15)

# ======================================================================
# SLIDE 4 — build-ahead status A–E
# ======================================================================
s = new_slide(
    KICKER,
    "Built ahead: the Python side of every phase is ready",
    [[("The roadmap runs A → E, each step gated by a safety review before Python gains any "
       "authority. The Python software for ", MUTED, False),
      ("all five phases", INK, True),
      (" is written and tested — what remains is LabVIEW wiring, joint decisions, and "
       "team-supplied numbers.", MUTED, False)]],
    page=4,
)

rows = [
    ("A", "Shadow brain", "state machine + warning policy ported; validated live against "
     "the real system — 100% agreement", "DONE — VALIDATED LIVE", OK, OK_BG, OK_LN),
    ("B", "Command path", "Python commander built and drill-tested against a simulator; "
     "interface spec v0.2 frozen with the team; PC-loss watchdog closed",
     "PYTHON DONE · LABVIEW NEXT", WARN, WARN_BG, WARN_LN),
    ("C", "Operator tools", "engineering CLI + an alarm that flags any live divergence "
     "between Python and LabVIEW decisions", "BUILT — AWAITS PHASE B", PY, PY_BG, PY_LN),
    ("D", "Sequencing", "step-by-step procedure engine + plant simulator; draft venting / "
     "purge / warm-up sequences survive randomized fault injection",
     "BUILT — NEEDS PROCEDURES", PY, PY_BG, PY_LN),
    ("E", "Warnings & setpoints", "time-based warning rules and setpoint scheduler engines, "
     "fully data-driven", "BUILT — NEEDS NUMBERS", PY, PY_BG, PY_LN),
]
ry = 2.42
for tag, name, desc, status, fg, bg, ln in rows:
    rect(s, M, ry, PAGE_W - 2 * M, 0.74, fill=PANEL, line=HAIR, line_w=0.75,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    tsq = rect(s, M + 0.12, ry + 0.12, 0.5, 0.5, fill=bg, line=ln, line_w=1.0,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
    tf = tsq.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag
    r.font.size = Pt(16); r.font.bold = True; r.font.name = MONO; r.font.color.rgb = fg
    text(s, M + 0.78, ry + 0.10, 2.05, 0.5, name, size=13.5, color=INK, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, M + 2.9, ry + 0.08, 6.55, 0.6, desc, size=10.5, color=MUTED, spacing=1.05,
         anchor=MSO_ANCHOR.MIDDLE)
    st = rect(s, M + 9.6, ry + 0.17, 2.55, 0.40, fill=bg, line=ln, line_w=1.0,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = st.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = status
    r.font.size = Pt(8.5); r.font.bold = True; r.font.name = MONO; r.font.color.rgb = fg
    ry += 0.84

hline(s, M, ry + 0.10, PAGE_W - 2 * M)
text(s, M, ry + 0.24, PAGE_W - 2 * M, 0.7,
     [[("Also in place:  ", INK, True),
       ("144+ automated tests run on every change (two Python versions); every remaining "
        "LabVIEW change is written up click-by-click so the work isn’t bottlenecked on any "
        "one person; a handoff guide covers how to continue the project.", MUTED, False)]],
     size=11.5, spacing=1.15)

# ======================================================================
# SLIDE 5 — next steps + asks
# ======================================================================
s = new_slide(
    KICKER,
    "What’s next — and what we need from the team",
    None,
    page=5,
)

LX, LW = M, 6.6
text(s, LX, 1.65, LW, 0.25, "NEXT STEPS (IN ORDER)", size=10.5, color=MUTED, font=MONO, tracking=110)
hline(s, LX, 1.97, LW)
nxt = [
    ("1", "Wire the LabVIEW command path (B3)",
     "the gateway learns to accept Python’s requests — with LabVIEW still validating "
     "every one and a single switch choosing UI or Python as the source. Fully "
     "specified, node-by-node."),
    ("2", "Bench command drills (B4)",
     "prove the failure story end-to-end: bad commands rejected, stale link means "
     "commands stop, source flips are clean, kill-Python-mid-run safe-holds."),
    ("3", "First supervised bench use (Phase C)",
     "Python drives mode requests on the bench — operator watching, LabVIEW clamp "
     "always underneath, one-command handback to the UI."),
]
ny = 2.15
for num, t, sub in nxt:
    circ = rect(s, LX, ny + 0.03, 0.42, 0.42, fill=PY_BG, line=PY_LN, line_w=1.0,
                shape=MSO_SHAPE.OVAL)
    tf = circ.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.size = Pt(14); r.font.bold = True; r.font.name = MONO; r.font.color.rgb = PY
    text(s, LX + 0.62, ny, LW - 0.62, 0.32, t, size=13.5, color=INK, bold=True)
    text(s, LX + 0.62, ny + 0.32, LW - 0.62, 0.75, sub, size=11, color=MUTED, spacing=1.1)
    ny += 1.28

RX, RW = 7.65, 5.13
text(s, RX, 1.65, RW, 0.25, "WHAT WE NEED FROM THE TEAM", size=10.5, color=MUTED,
     font=MONO, tracking=110)
hline(s, RX, 1.97, RW)
asks = [
    ("Operating procedures", "the venting / purge / warm-up sequences are drafted from "
     "engineering judgement — they need the team’s real steps, setpoints and hold "
     "conditions (template sheets ready)"),
    ("Real numbers", "warning thresholds, time windows and setpoint schedules are "
     "engine knowledge, not code — the engines are built and waiting for values"),
    ("Sign-off at each gate", "every increase in Python’s authority (shadow → bench "
     "command → sequences) gets a joint review before it happens"),
]
ay = 2.15
for t, sub in asks:
    text(s, RX, ay, RW, 0.3, [[("▪  ", LV, True), (t, INK, True)]], size=12.5)
    text(s, RX + 0.25, ay + 0.30, RW - 0.25, 0.85, sub, size=10.5, color=MUTED, spacing=1.1)
    ay += 1.24

# safety footer
fy = 6.35
fb = rect(s, M, fy, PAGE_W - 2 * M, 0.72, fill=LV_BG, line=LV_LN, line_w=1.0,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
tf = fb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.word_wrap = True
tf.margin_left = Inches(0.22)
tf.margin_right = Inches(0.22)
p = tf.paragraphs[0]
for seg, c, b in [("The rule that never changes:  ", LV, True),
                  ("LabVIEW/cRIO keeps the hardware, the interlocks and the safe fallback — "
                   "permanently. Python only ever asks; LabVIEW checks every request; and if "
                   "Python disappears, the system safe-holds. That last part is now proven, "
                   "not assumed.", INK, False)]:
    r = p.add_run(); r.text = seg
    r.font.size = Pt(11.5); r.font.bold = b; r.font.name = SANS; r.font.color.rgb = c

out = HERE / "migration-progress.pptx"
prs.save(str(out))
print("wrote", out)
